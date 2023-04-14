import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches

#Creation of a Powerpoint Presentation
pr1=Presentation()



##Creation of First Slide in the Presentation
#The number 0 refers to the layout used in powerpoint Presentation for adding the content resp.
slide1_register = pr1.slide_layouts[0]
#adding this slide into the presentation
slide1 = pr1.slides.add_slide(slide1_register) 
#Succesfully Created Presentation with slide1.

#Adding Content to slide 1
#layouts have specific placeholders that hold the content in the slide
#for eg:layout1 have two placeholders title and subtitle and layout2 has Two  placeholder, one for title and the other for image.
title = slide1.shapes.title
title_text = input("Enter title of the Project:")
title.text = title_text
#Putting data into subtitle
subtitle = slide1.placeholders[1]
subtitle_text = input("Enter Project Members:")
subtitle.text = subtitle_text



##Creation of Second Slide in the Presentation
# Creating a bullet point slide
slide2_register = pr1.slide_layouts[1]
slide2 = pr1.slides.add_slide(slide2_register)

title2 = slide2.shapes.title
title2.text = "About"

bullet_point_box = slide2.shapes

bullet_points_lvl1 = bullet_point_box.placeholders[1]
bullet_points_lvl1.text = "Computer"

bullet_points_lvl2 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl2.text = "is"
bullet_points_lvl2.level=1

bullet_points_lvl3 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl3.text = "my"
bullet_points_lvl3.level=2

bullet_points_lvl4 = bullet_points_lvl1.text_frame.add_paragraph()
bullet_points_lvl4.text = "Favroite tool"
bullet_points_lvl4.level=3




#creation of slide 3
#creating a slide that can hold a title and a picture
slide3_register = pr1.slide_layouts[5]
slide3 = pr1.slides.add_slide(slide3_register)

title3 = slide3.shapes.title
title3.text = "Picture!!"

#Adding images
img1 = "computer.png"
from_left = Inches(2)
from_top = Inches (2)
add_picture = slide3.shapes.add_picture(img1,from_left,from_top)

#To save our presentation
pr1.save("PPT.pptx")
